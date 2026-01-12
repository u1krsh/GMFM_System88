"""
Generate a premium, modern PowerPoint presentation for MotorMeasure Pro
With gradients, modern design, and comprehensive content for first review
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Get script directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Modern color palette - Premium medical tech feel
DEEP_TEAL = RGBColor(0x0D, 0x94, 0x88)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_CYAN = RGBColor(0xCC, 0xFB, 0xF1)

# Accent colors
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

# ============== TEAM DETAILS - EDIT THESE ==============
PRODUCT_NAME = "MotorMeasure Pro"
PRODUCT_TAGLINE = "Next-Generation Clinical Assessment Platform"
PRODUCT_SUBTITLE = "Transforming pediatric motor function evaluation with intelligent digital tools"

# Team members - ADD YOUR DETAILS HERE
TEAM_MEMBERS = [
    {"name": "Team Member 1", "roll": "XXX123", "role": "Developer"},
    {"name": "Team Member 2", "roll": "XXX124", "role": "Developer"},
    {"name": "Team Member 3", "roll": "XXX125", "role": "Developer"},
    {"name": "Team Member 4", "roll": "XXX126", "role": "Developer"},
]
GUIDE_NAME = "Dr. Guide Name"
DEPARTMENT = "Department of Computer Science & Engineering"
INSTITUTION = "Your Institution Name"
# =====================================================


def add_gradient_bg(slide, color1, color2):
    """Add a subtle gradient-like effect using multiple rectangles"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1
    bg.line.fill.background()
    
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4), prs.slide_width, Inches(3.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color2
    overlay.fill.fore_color.brightness = 0.05
    overlay.line.fill.background()


def add_hero_slide_with_team(prs, title, subtitle, tagline=""):
    """Create a stunning hero/title slide with team details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    # Decorative accent circles (top right)
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1.5), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = DEEP_TEAL
    circle1.fill.fore_color.brightness = -0.3
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-0.5), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = EMERALD
    circle2.fill.fore_color.brightness = -0.4
    circle2.line.fill.background()
    
    # Decorative accent (bottom left)
    circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5.5), Inches(4), Inches(4))
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = VIOLET
    circle3.fill.fore_color.brightness = -0.5
    circle3.line.fill.background()
    
    # Glowing accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(1.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DEEP_TEAL
    accent.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(10), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle with accent color
    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(10), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_CYAN
    
    # Department and Institution
    dept_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(8), Inches(0.5))
    tf = dept_box.text_frame
    p = tf.paragraphs[0]
    p.text = DEPARTMENT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p = tf.add_paragraph()
    p.text = INSTITUTION
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    # Team members in cards at bottom
    team_start_x = Inches(0.5)
    team_y = Inches(4.8)
    card_width = Inches(2.9)
    
    for i, member in enumerate(TEAM_MEMBERS[:4]):
        x = team_start_x + i * (card_width + Inches(0.3))
        
        # Member card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, team_y, card_width, Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_800
        card.line.color.rgb = SLATE_700
        card.line.width = Pt(1)
        
        # Member info
        name_box = slide.shapes.add_textbox(x + Inches(0.15), team_y + Inches(0.2), card_width - Inches(0.3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = member["name"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p = tf.add_paragraph()
        p.text = member["roll"]
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_CYAN
        
        p = tf.add_paragraph()
        p.text = member.get("role", "Team Member")
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    # Guide info at bottom right
    guide_box = slide.shapes.add_textbox(Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.8))
    tf = guide_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Guide: {GUIDE_NAME}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_modern_content_slide(prs, title, bullets, accent=DEEP_TEAL):
    """Modern content slide with glassmorphism feel"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    dec = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(6), Inches(6))
    dec.fill.solid()
    dec.fill.fore_color.rgb = accent
    dec.fill.fore_color.brightness = -0.6
    dec.line.fill.background()
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(0.12), Inches(0.7))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    start_y = Inches(1.8)
    for i, (icon, text) in enumerate(bullets):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.2)
        y = start_y + row * Inches(1.4)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.9), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_800
        card.line.color.rgb = SLATE_700
        card.line.width = Pt(1)
        
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.7))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = accent
        icon_bg.fill.fore_color.brightness = -0.2
        icon_bg.line.fill.background()
        
        icon_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.35), Inches(0.7), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        
        text_box = slide.shapes.add_textbox(x + Inches(1.1), y + Inches(0.3), Inches(4.5), Inches(0.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
    
    return slide


def add_section_divider(prs, title, subtitle="", color=DEEP_TEAL):
    """Bold section divider with large typography"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    
    pattern = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8))
    pattern.fill.solid()
    pattern.fill.fore_color.rgb = WHITE
    pattern.fill.fore_color.brightness = -0.9
    pattern.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_feature_showcase(prs, title, features):
    """Premium feature showcase with icons and descriptions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    colors = [DEEP_TEAL, EMERALD, SKY, VIOLET, AMBER, ROSE]
    
    for i, (feat_title, feat_desc) in enumerate(features[:6]):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(2.9)
        color = colors[i % len(colors)]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_800
        card.line.fill.background()
        
        top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Inches(0.15))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = color
        top_accent.line.fill.background()
        
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.5), Inches(0.8), Inches(0.8))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = color
        icon_bg.fill.fore_color.brightness = -0.2
        icon_bg.line.fill.background()
        
        num_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6), Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        feat_title_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.4), Inches(3.3), Inches(0.5))
        tf = feat_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.85), Inches(3.3), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feat_desc
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    
    return slide


def add_two_column_premium(prs, title, left, right):
    """Premium two-column layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(5.6))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = SLATE_800
    left_card.line.fill.background()
    
    left_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(0.8))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = DEEP_TEAL
    left_header.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left["title"]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    left_content = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.4), Inches(4.4))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left["items"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "> " + item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6), Inches(5.6))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = SLATE_800
    right_card.line.fill.background()
    
    right_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.4), Inches(6), Inches(0.8))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = EMERALD
    right_header.line.fill.background()
    
    right_title_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right["title"]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    right_content = slide.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.4), Inches(4.4))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right["items"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "> " + item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
    
    return slide


def add_stats_slide(prs, title, stats):
    """Big stats/metrics showcase"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    colors = [DEEP_TEAL, EMERALD, VIOLET, SKY]
    card_width = Inches(2.9)
    start_x = Inches(0.6)
    
    for i, (number, label) in enumerate(stats[:4]):
        x = start_x + i * (card_width + Inches(0.3))
        color = colors[i]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), card_width, Inches(3))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_800
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        num_box = slide.shapes.add_textbox(x, Inches(2.4), card_width, Inches(1.2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.7), card_width - Inches(0.4), Inches(1))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_workflow_slide(prs, title, steps):
    """Horizontal workflow/process slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    colors = [DEEP_TEAL, SKY, EMERALD, VIOLET, AMBER, ROSE]
    step_width = Inches(2)
    start_x = Inches(0.5)
    
    for i, (step_title, step_desc) in enumerate(steps[:6]):
        x = start_x + i * Inches(2.1)
        color = colors[i]
        
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), Inches(1.8), Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(x + Inches(0.5), Inches(2), Inches(1), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.6), Inches(2.15), Inches(0.5), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SLATE_700
            arrow.line.fill.background()
        
        title_box = slide.shapes.add_textbox(x, Inches(3), step_width, Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x, Inches(3.5), step_width, Inches(2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_text_slide(prs, title, paragraphs, accent=DEEP_TEAL):
    """Slide with title and multiple paragraphs of text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    # Decorative element
    dec = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(-1), Inches(4), Inches(4))
    dec.fill.solid()
    dec.fill.fore_color.rgb = accent
    dec.fill.fore_color.brightness = -0.6
    dec.line.fill.background()
    
    # Title bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(0.12), Inches(0.7))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(18)
    
    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Slide with an embedded image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add image if exists
    if os.path.exists(image_path):
        # Calculate image size to fit slide
        img_left = Inches(0.5)
        img_top = Inches(1.2)
        img_width = Inches(12.3)
        img_height = Inches(5.5)
        
        slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
    else:
        # Placeholder if image not found
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = SLATE_800
        placeholder.line.color.rgb = SLATE_700
        
        text_box = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(5), Inches(1))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {os.path.basename(image_path)}]"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_closing_slide(prs, title, subtitle, cta=""):
    """Premium closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE_900
    bg.line.fill.background()
    
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(7), Inches(7))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = DEEP_TEAL
    circle1.fill.fore_color.brightness = -0.5
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(6), Inches(6))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = VIOLET
    circle2.fill.fore_color.brightness = -0.6
    circle2.line.fill.background()
    
    circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = EMERALD
    circle3.fill.fore_color.brightness = -0.5
    circle3.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    if cta:
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(5.2), Inches(3.3), Inches(0.8))
        btn.fill.solid()
        btn.fill.fore_color.rgb = DEEP_TEAL
        btn.line.fill.background()
        
        btn_text = slide.shapes.add_textbox(Inches(5), Inches(5.35), Inches(3.3), Inches(0.6))
        tf = btn_text.text_frame
        p = tf.paragraphs[0]
        p.text = cta
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


# ============ CREATE THE PRESENTATION ============

# Slide 1: Hero with Team
add_hero_slide_with_team(prs, 
    PRODUCT_NAME,
    PRODUCT_TAGLINE,
    PRODUCT_SUBTITLE)

# Slide 2: About Cerebral Palsy
add_text_slide(prs, "What is Cerebral Palsy (CP)?", [
    "Cerebral Palsy is a group of permanent movement disorders that appear in early childhood, affecting muscle tone, movement, and motor skills.",
    "It is caused by abnormal brain development or damage to the developing brain, typically before birth or during the first few years of life.",
    "CP is the most common motor disability in childhood, affecting approximately 2-3 children per 1,000 live births worldwide.",
    "Children with CP often experience difficulties with posture, balance, and coordination, making standardized motor function assessment crucial for treatment planning.",
    "Early intervention and regular assessment are key to improving outcomes. This is where standardized motor assessment tools become essential."
], ROSE)

# Slide 3: About Spasticity
add_text_slide(prs, "Understanding Spasticity", [
    "Spasticity is a condition characterized by abnormal muscle tightness (increased muscle tone) due to prolonged muscle contraction. It is one of the most common symptoms in Cerebral Palsy.",
    "In CP patients, spasticity occurs because of damage to the parts of the brain that control voluntary movement, resulting in overactive reflexes and stiff muscles.",
    "Spasticity significantly impacts motor function by limiting range of motion, causing pain, and making everyday movements difficult.",
    "Regular assessment of motor function helps clinicians track the effectiveness of spasticity management interventions including physical therapy, medications, and botulinum toxin injections.",
    "Our platform enables systematic tracking of motor abilities to measure treatment outcomes over time."
], VIOLET)

# Slide 4: Why Motor Assessment Matters
add_modern_content_slide(prs, "Why Standardized Motor Assessment?", [
    ("1", "Objective Measurement: Provides consistent, reproducible scores across different evaluators"),
    ("2", "Treatment Planning: Helps therapists design targeted intervention programs"),
    ("3", "Progress Tracking: Enables monitoring of improvement or decline over time"),
    ("4", "Research: Facilitates clinical research with standardized outcome measures"),
    ("5", "Communication: Creates common language among healthcare providers"),
    ("6", "Documentation: Provides legal and medical documentation of patient status"),
])

# Slide 5: Section - The Problem
add_section_divider(prs, "The Problem", "Understanding the challenges faced today", ROSE)

# Slide 6: Target Audience
add_two_column_premium(prs, "Who Needs This Solution?",
    {
        "title": "Primary Users",
        "items": [
            "Physiotherapists in clinics & hospitals",
            "Pediatric neurologists",
            "Occupational therapists",
            "Rehabilitation specialists",
            "Special education professionals",
            "Researchers in motor disorders"
        ]
    },
    {
        "title": "Settings",
        "items": [
            "Pediatric rehabilitation centers",
            "Multi-specialty hospitals",
            "Special needs schools",
            "Home healthcare services",
            "Rural and urban clinics",
            "Academic research institutions"
        ]
    }
)

# Slide 7: Problem Statement with image
add_image_slide(prs, 
    "Problem Space Analysis",
    os.path.join(SCRIPT_DIR, "problem_diagram.png"),
    "Current challenges in motor function assessment workflow")

# Slide 8: Existing System Analysis
add_two_column_premium(prs, "Existing System Analysis",
    {
        "title": "Current Paper-Based Method",
        "items": [
            "88-item paper scoresheets required",
            "Manual calculation of domain percentages",
            "High risk of arithmetic errors",
            "Physical storage of records needed",
            "Difficult to track progress over time",
            "Time-consuming data entry for research"
        ]
    },
    {
        "title": "Problems with Current Approach",
        "items": [
            "No standardized digital workflow",
            "Inconsistent scoring between sessions",
            "Limited ability to compare results",
            "No automated report generation",
            "Data security concerns with paper records",
            "Inefficient for large patient volumes"
        ]
    }
)

# Slide 9: Market Size
add_stats_slide(prs, "Market Opportunity", [
    ("17M+", "CP patients worldwide need regular assessment"),
    ("2-3", "Per 1000 live births affected annually"),
    ("$8.5B", "Global rehabilitation market (growing 6% YoY)"),
    ("0", "Affordable Indian solutions currently available"),
])

# Slide 10: Section - Our Solution
add_section_divider(prs, "Our Solution", f"{PRODUCT_NAME} - Digital Assessment Platform", DEEP_TEAL)

# Slide 11: What is our Assessment
add_modern_content_slide(prs, "About the Assessment Method", [
    ("1", "Standardized 88-item assessment across 5 developmental domains"),
    ("2", "Designed specifically for children with cerebral palsy and motor disorders"),
    ("3", "Measures motor abilities: lying, sitting, crawling, standing, walking"),
    ("4", "Scoring system: 0 (cannot initiate) to 3 (completes independently)"),
    ("5", "Used by therapists and clinicians worldwide as the gold standard"),
    ("6", "Our digital platform faithfully implements this proven methodology"),
])

# Slide 12: The 5 Domains
add_feature_showcase(prs, "Assessment Domains", [
    ("A: Lying & Rolling", "Basic floor mobility and positional changes - 17 items"),
    ("B: Sitting", "Seated balance, transitions and functional sitting - 20 items"),
    ("C: Crawling & Kneeling", "Pre-walking locomotion and floor mobility - 14 items"),
    ("D: Standing", "Upright posture, balance and standing tasks - 13 items"),
    ("E: Walking & Running", "Advanced locomotion, jumping and stairs - 24 items"),
    ("Total: 88 Items", "Complete motor function profile across all domains"),
])

# Slide 13: Section - Architecture
add_section_divider(prs, "System Architecture", "Technical excellence in design", VIOLET)

# Slide 14: Architecture Diagram
add_image_slide(prs,
    "System Architecture Overview",
    os.path.join(SCRIPT_DIR, "architecture_diagram.png"),
    "Clean, layered architecture design")

# Slide 15: Architecture Details
add_two_column_premium(prs, "Architectural Components",
    {
        "title": "Application Layer",
        "items": [
            "Flet Python UI Framework",
            "View-based navigation system", 
            "Route parameters for data passing",
            "Async-safe database operations",
            "Thread-based session timer",
            "Event-driven reactive UI"
        ]
    },
    {
        "title": "Data Layer",
        "items": [
            "SQLite with encrypted fields",
            "Pydantic models for validation",
            "Repository pattern for data access",
            "JSON catalog of 88 test items",
            "Embedded exercise instructions",
            "PDF export service"
        ]
    }
)

# Slide 16: Security Architecture
add_modern_content_slide(prs, "Enterprise-Grade Security", [
    ("A", "Field-level AES encryption for student names and identifiers"),
    ("B", "Keys stored securely in OS keyring (Keychain/Credential Manager)"),
    ("C", "Local SQLite database - zero cloud data transmission"),
    ("D", "Pydantic validation prevents malformed data entry"),
    ("E", "All data stays on device - HIPAA-friendly approach"),
    ("F", "Fallback encryption from environment variables"),
], ROSE)

# Slide 17: Workflow Diagram
add_image_slide(prs,
    "Clinical Workflow",
    os.path.join(SCRIPT_DIR, "workflow_diagram.png"),
    "Streamlined end-to-end assessment process")

# Slide 18: Workflow Steps
add_workflow_slide(prs, "Assessment Process", [
    ("Add Patient", "Enter name, DOB and identifier"),
    ("Select Scale", "Choose 66 or 88 item scale"),
    ("Navigate", "Tabs organize domains A-E"),
    ("Score Items", "Tap 0-3 or NT for each"),
    ("View Help", "Tap item for instructions"),
    ("Save Results", "Auto-calculated with notes"),
])

# Slide 19: Section - Features
add_section_divider(prs, "Application Features", f"Discover what makes {PRODUCT_NAME} powerful", EMERALD)

# Slide 20: Dashboard Features
add_feature_showcase(prs, "Intelligent Dashboard", [
    ("Patient Management", "Add, edit, and organize patient profiles with secure storage"),
    ("Live Statistics", "Real-time counts of patients, sessions, and average scores"),
    ("Recent Activity", "Quick access to the last 3 assessment sessions"),
    ("Smart Search", "Instant search with real-time text highlighting"),
    ("Quick Actions", "One-tap access to new patient or DOCX import"),
    ("Dynamic Greeting", "Time-based personalized welcome messages"),
])

# Slide 21: Scoring Features
add_feature_showcase(prs, "Professional Scoring Interface", [
    ("Digital Scoresheet", "Mirrors the official paper form perfectly"),
    ("One-Tap Scoring", "Touch 0, 1, 2, 3 or NT with instant visual feedback"),
    ("Live Progress", "Real-time progress bar tracks completion percentage"),
    ("Session Timer", "Automatic duration tracking for documentation"),
    ("Bulk Actions", "Score all items as 0 or 3, or clear all with one tap"),
    ("Jump to Unscored", "Instantly navigate to the next missing item"),
])

# Slide 22: Cross Platform
add_two_column_premium(prs, "Cross-Platform Design",
    {
        "title": "Platform Support",
        "items": [
            "Android smartphones and tablets",
            "iOS devices (iPhone & iPad)",
            "Windows desktop application",
            "macOS desktop support",
            "Single Python/Flet codebase",
            "Native look and feel on each platform"
        ]
    },
    {
        "title": "User Experience",
        "items": [
            "Light and Dark theme support",
            "Haptic feedback on interactions",
            "Smooth gesture-based navigation",
            "Safe area handling for notches",
            "Responsive adaptive layouts",
            "Celebration animations on completion"
        ]
    }
)

# Slide 23: Services
add_feature_showcase(prs, "Integrated Services", [
    ("Scoring Engine", "Accurate 66/88 calculations with domain analysis"),
    ("Haptics Service", "Native feedback for tap, select, and success events"),
    ("Instructions DB", "Static JSON with detailed exercise guidance"),
    ("DOCX Import", "Parse existing assessments from Word documents"),
    ("Chart Service", "Generate progress trend visualizations"),
    ("PDF Export", "Official format scoresheets with marked scores"),
])

# Slide 24: Why Choose Us
add_two_column_premium(prs, f"Why Choose {PRODUCT_NAME}?",
    {
        "title": "For Clinicians",
        "items": [
            "Replaces paper scoresheets entirely",
            "Instant automatic score calculations",
            "Eliminates manual math errors",
            "Built-in exercise instructions",
            "Track patient progress over time",
            "Secure compliant data storage"
        ]
    },
    {
        "title": "For Organizations",
        "items": [
            "Offline-first - works anywhere",
            "No subscription fees required",
            "HIPAA-friendly encryption",
            "Cross-platform deployment",
            "Open extensible architecture",
            "Professional PDF documentation"
        ]
    }
)

# Slide 25: Roadmap
add_modern_content_slide(prs, "Future Development Roadmap", [
    ("1", "Remote sync adapter for multi-device access and backup"),
    ("2", "Role-based access control for clinical teams"),
    ("3", "Integration with Electronic Health Records (EHR)"),
    ("4", "Multi-language support for global deployment"),
    ("5", "Advanced analytics and reporting dashboard"),
    ("6", "Cloud backup and secure data restore options"),
], AMBER)

# Slide 26: Closing
add_closing_slide(prs,
    "Thank You",
    f"{PRODUCT_NAME} - Transforming Motor Function Assessment",
    "Questions?")

# Save
output_path = os.path.join(SCRIPT_DIR, "MotorMeasure_Pro_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Design: Modern dark theme with gradient accents")
print("\nNOTE: Edit TEAM_MEMBERS, GUIDE_NAME, DEPARTMENT, and INSTITUTION")
print("      at the top of this file to customize the title slide.")
